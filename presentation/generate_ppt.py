import sys
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: The 'python-pptx' library is required to run this script.")
    print("Please install it using: pip install python-pptx")
    sys.exit(1)

def create_presentation():
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "RouteMaster Live"
    subtitle.text = "Real-Time Supply Chain Command Center\nProject Overview"

    # Slide 2: The Challenge
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Challenge"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Modern supply chains face critical visibility issues:"
    p = tf.add_paragraph()
    p.text = "Lack of real-time tracking for assets and fleets."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Disconnected systems leading to data silos."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Inefficient route planning and delays."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Difficulty in monitoring driver performance and safety."
    p.level = 1

    # Slide 3: Our Solution
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Our Solution: RouteMaster Live"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "A comprehensive command center for logistics management:"
    p = tf.add_paragraph()
    p.text = "Real-time asset tracking on interactive maps."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Instant alerts and notifications via WebSocket."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Data-driven insights for route optimization."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Seamless integration with existing ERP systems."
    p.level = 1

    # Slide 4: Key Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Features"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Robust functionality designed for scale:"
    p = tf.add_paragraph()
    p.text = "Live Map Visualization (Leaflet & OpenStreetMap)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Real-time Data Feeds (WebSocket/STOMP)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Advanced Analytics Dashboard (Chart.js)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Secure Role-Based Access Control (JWT)."
    p.level = 1

    # Slide 5: Technical Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technical Architecture"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Built on a modern, scalable tech stack:"
    p = tf.add_paragraph()
    p.text = "Frontend: Vue.js 3, Vite, Pinia, Tailwind CSS"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend: Spring Boot 3.2, Spring Security"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Database: MongoDB (NoSQL), H2 (Batch Metadata)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Real-time: WebSocket (STOMP Protocol)"
    p.level = 1

    # Slide 6: Backend Deep Dive
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Backend Overview"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Core backend components:"
    p = tf.add_paragraph()
    p.text = "Spring Boot: Rapid application development framework."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "MongoDB: Flexible document storage for complex data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Spring Batch: Efficient processing of large datasets."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Security: Stateless JWT authentication & authorization."
    p.level = 1

    # Slide 7: Frontend Deep Dive
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Frontend Overview"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "User-centric interface design:"
    p = tf.add_paragraph()
    p.text = "Vue 3 Composition API: Modular and reactive code base."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Pinia: Centralized state management."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Visuals: Integrated Chart.js graphs and Leaflet maps."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Responsiveness: Optimized for desktop and mobile devices."
    p.level = 1

    # Slide 8: Data Flow Roadmap
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "System Data Flow"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "End-to-end data journey:"
    p = tf.add_paragraph()
    p.text = "1. Data Ingestion: IoT sensors/API input -> Backend."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Processing: Spring Batch normalization & validation."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Storage: Persisted securely in MongoDB."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Delivery: Real-time push via WebSocket to Frontend."
    p.level = 1

    # Slide 9: Future Enhancements
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Future Roadmap"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Planned upgrades for upcoming releases:"
    p = tf.add_paragraph()
    p.text = "AI-driven predictive analytics for route optimization."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Dedicated mobile application for drivers."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Blockchain integration for transparent audit trails."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Integration with third-party logistics providers (3PL)."
    p.level = 1

    # Slide 10: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "RouteMaster Live is positioned to transform logistics:"
    p = tf.add_paragraph()
    p.text = "Scalable architecture ready for enterprise deployment."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Secure, real-time feedback loops."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Modern tech stack ensuring long-term maintainability."
    p.level = 1

    output_file = "RouteMaster_Live_Presentation.pptx"
    prs.save(output_file)
    print(f"Successfully created presentation: {output_file}")

if __name__ == "__main__":
    create_presentation()
